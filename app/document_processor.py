"""
Document Processor
Handles extraction of text and images from various file formats (PDF, DOCX, TXT)
"""
import os
import re
from typing import List, Optional, Tuple

import pdfplumber  # type: ignore
from docx import Document  # type: ignore

# Gemini / OpenAI vision APIs reject EMF/WMF (common in Word paste-from-Office).
_VISION_UNSUPPORTED_MIMES = frozenset(
    {
        "image/x-emf",
        "image/emf",
        "image/x-wmf",
        "image/wmf",
        "image/x-ms-bmp",  # sometimes mis-tagged vector blobs
    }
)
_VISION_PREFERRED_MIMES = (
    "image/png",
    "image/jpeg",
    "image/jpg",
    "image/webp",
    "image/gif",
    "image/bmp",
    "image/tiff",
)


def _normalize_vision_mime(mime_type: str) -> str:
    return (mime_type or "image/png").lower().split(";")[0].strip()


def is_vision_supported_mime(mime_type: str) -> bool:
    mime = _normalize_vision_mime(mime_type)
    if mime in _VISION_UNSUPPORTED_MIMES:
        return False
    return mime.startswith("image/")


def filter_vision_images(
    images: List[Tuple[bytes, str]],
) -> List[Tuple[bytes, str]]:
    """Drop EMF/WMF and prefer raster formats for vision API submission."""
    supported: List[Tuple[bytes, str]] = []
    for image_bytes, mime_type in images:
        if not image_bytes or len(image_bytes) < 2048:
            continue
        mime = _normalize_vision_mime(mime_type)
        if not is_vision_supported_mime(mime):
            continue
        supported.append((image_bytes, mime))

    def _rank(item: Tuple[bytes, str]) -> tuple:
        mime = item[1]
        try:
            pref = _VISION_PREFERRED_MIMES.index(mime)
        except ValueError:
            pref = len(_VISION_PREFERRED_MIMES)
        return (pref, -len(item[0]))

    supported.sort(key=_rank)
    return supported


def extract_text_from_file(file_path: str) -> str:
    """Extract text from file - convenience wrapper around DocumentProcessor.extract_text"""
    return DocumentProcessor.extract_text(file_path)


def extract_student_name_from_file(file_path: str) -> Optional[str]:
    """Extract student name from the first 1-2 pages of a Word/PDF document.

    Looks for patterns like:
      الطالب: اسم الطالب
      الاسم: اسم الطالب
      اسم الطالب: اسم الطالب
      Student Name: Name
      Learner Name: Name
      Name: Name
    """
    file_ext = os.path.splitext(file_path)[1].lower()

    try:
        if file_ext in ('.docx', '.doc'):
            return _extract_name_from_docx(file_path)
        elif file_ext == '.pdf':
            return _extract_name_from_pdf(file_path)
    except Exception as e:
        print(f"⚠️ Could not extract student name from {file_path}: {e}")

    return None


# Regex patterns to match student name lines (Arabic + English)
_NAME_PATTERNS = [
    # Arabic patterns
    re.compile(r'(?:اسم\s*الطالب[ة]?|الطالب[ة]?|الاسم)\s*[:：]\s*(.+)', re.UNICODE),
    # English patterns
    re.compile(r'(?:student\s*name|learner\s*name|name)\s*[:：]\s*(.+)', re.IGNORECASE),
]


def _clean_extracted_name(raw: str) -> Optional[str]:
    """Clean and validate an extracted name."""
    name = raw.strip().strip('.-:،,')
    # Remove trailing noise (numbers, codes, dates)
    name = re.sub(r'\s*\d{4,}.*$', '', name)
    name = re.sub(r'\s*TF\d+.*$', '', name, flags=re.IGNORECASE)
    name = name.strip()
    # Must be at least 3 chars and mostly letters (Arabic or Latin)
    if len(name) < 3:
        return None
    letter_count = sum(1 for c in name if c.isalpha())
    if letter_count < len(name) * 0.5:
        return None
    return name


def _extract_name_from_docx(file_path: str) -> Optional[str]:
    """Extract student name from first ~40 paragraphs of a DOCX."""
    doc = Document(file_path)
    for para in doc.paragraphs[:40]:
        text = para.text.strip()
        if not text:
            continue
        for pattern in _NAME_PATTERNS:
            m = pattern.search(text)
            if m:
                name = _clean_extracted_name(m.group(1))
                if name:
                    return name
    return None


def _extract_name_from_pdf(file_path: str) -> Optional[str]:
    """Extract student name from first 2 pages of a PDF."""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages[:2]:
            page_text = page.extract_text() or ""
            for line in page_text.split('\n'):
                line = line.strip()
                if not line:
                    continue
                for pattern in _NAME_PATTERNS:
                    m = pattern.search(line)
                    if m:
                        name = _clean_extracted_name(m.group(1))
                        if name:
                            return name
    return None


class DocumentProcessor:
    """
    Handles file processing and text extraction
    """

    @staticmethod
    def extract_text(file_path: str, max_chars: Optional[int] = None) -> str:
        """
        Extract text from PDF, DOCX, or TXT files

        Args:
            file_path: Absolute path to the file
            max_chars: Optional cap — only when callers pass it; grading uses full extract

        Returns:
            Extracted text content
        """
        from pathlib import Path as _Path

        resolved = _Path(file_path).resolve()
        file_path = str(resolved)
        if not resolved.is_file():
            raise FileNotFoundError(f"File not found: {file_path}")

        file_ext = resolved.suffix.lower()

        try:
            # Extended code extensions covering all BTEC IT (Jordan) units:
            #   Programming, OOP, Web, Mobile (Flutter/Android/iOS),
            #   Game Dev (Unity/Unreal/Scratch), AI, Big Data, Networking.
            _CODE_EXTENSIONS = (
                # General programming
                '.py', '.java', '.cs', '.cpp', '.c', '.h', '.hpp',
                '.js', '.ts', '.jsx', '.tsx', '.html', '.css', '.scss', '.less',
                '.rb', '.go', '.rs', '.php', '.swift', '.kt', '.scala', '.r',
                '.sql', '.sh', '.bat', '.ps1',
                # Configs / structured text
                '.xml', '.yaml', '.yml', '.ini', '.cfg', '.conf', '.log',
                '.csv', '.tsv', '.toml', '.env',
                # Game Dev source / project files
                '.gml',                          # GameMaker
                '.gmx',                          # GameMaker old
                '.uc',                           # Unreal script
                '.gd',                           # Godot
                # Mobile (Flutter / Dart / Android / iOS)
                '.dart', '.pubspec',             # Flutter
                '.gradle', '.kts', '.kotlin',    # Android (Kotlin)
                '.plist', '.storyboard', '.xib', # iOS
                # Data Science / AI / Notebooks (treated as text — JSON inside)
                '.ipynb', '.rmd', '.qmd',
                # Database schema / dump files
                '.sqlfile', '.ddl', '.dml',
                # Web Frameworks
                '.vue', '.svelte', '.razor', '.cshtml',
                # Misc
                '.md', '.markdown', '.rst', '.tex',
            )
            _IMAGE_EXTENSIONS = (
                '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.tiff',
                '.ico', '.heic', '.heif', '.psd', '.ai', '.eps',
            )
            _VIDEO_EXTENSIONS = (
                '.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm',
                '.m4v', '.3gp', '.mpeg', '.mpg',
            )
            _AUDIO_EXTENSIONS = (
                '.mp3', '.wav', '.ogg', '.m4a', '.flac', '.aac', '.wma',
            )
            _DATABASE_EXTENSIONS = (
                '.db', '.sqlite', '.sqlite3', '.mdb', '.accdb', '.dbf',
            )
            _NETWORKING_EXTENSIONS = (
                # Cisco Packet Tracer, Wireshark, Network designers
                '.pkt', '.pka', '.pcap', '.pcapng', '.cap',
            )
            _GAME_PROJECT_EXTENSIONS = (
                # Unity / Unreal / GameMaker / Scratch / Roblox
                '.unity', '.prefab', '.scene', '.unityproj', '.uasset',
                '.umap', '.gms2', '.yyp', '.gmproj', '.fla', '.swf',
                '.gma', '.rbxl', '.rbxmx',
            )
            _DESIGN_EXTENSIONS = (
                # Design / 3D / CAD
                '.fig', '.sketch', '.xd', '.afdesign', '.afphoto',
                '.blend', '.fbx', '.obj', '.stl', '.dwg', '.dxf',
            )

            if file_ext == '.pdf':
                return DocumentProcessor._extract_from_pdf(file_path)
            elif file_ext in ['.docx', '.doc']:
                return DocumentProcessor._extract_from_docx(file_path, max_chars=max_chars)
            elif file_ext in ['.xlsx', '.xls', '.xlsm', '.ods']:
                return DocumentProcessor._extract_from_excel(file_path)
            elif file_ext in ['.pptx', '.ppt', '.odp']:
                return DocumentProcessor._extract_from_pptx(file_path)
            elif file_ext == '.svg':
                return DocumentProcessor._extract_from_svg(file_path)
            elif file_ext in _VIDEO_EXTENSIONS:
                return DocumentProcessor._extract_video_metadata(file_path)
            elif file_ext in _AUDIO_EXTENSIONS:
                return DocumentProcessor._describe_binary_file(
                    file_path, "Audio recording",
                    "تسجيل صوتي مقدم كدليل")
            elif file_ext == '.sb3':
                return DocumentProcessor._extract_from_scratch(file_path)
            elif file_ext in _NETWORKING_EXTENSIONS:
                return DocumentProcessor._describe_binary_file(
                    file_path, "Cisco Packet Tracer / Network capture",
                    "ملف محاكاة شبكات Cisco Packet Tracer / تسجيل شبكة - "
                    "يُعتبر دليلاً عملياً على تنفيذ تصميم الشبكة")
            elif file_ext in _GAME_PROJECT_EXTENSIONS:
                return DocumentProcessor._describe_binary_file(
                    file_path, "Game project file",
                    "ملف مشروع لعبة (Unity/Unreal/GameMaker/Scratch) - "
                    "يُعتبر دليلاً على تطوير اللعبة")
            elif file_ext in _DATABASE_EXTENSIONS:
                return DocumentProcessor._describe_binary_file(
                    file_path, "Database file",
                    "ملف قاعدة بيانات - يُعتبر دليلاً على تصميم وتنفيذ قاعدة البيانات")
            elif file_ext in _DESIGN_EXTENSIONS:
                return DocumentProcessor._describe_binary_file(
                    file_path, "Design / CAD file",
                    "ملف تصميم (Figma/Sketch/CAD/3D) - يُعتبر دليلاً مرئياً")
            elif file_ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
                return DocumentProcessor._describe_binary_file(
                    file_path, "Archive (ZIP/RAR/7z)",
                    "أرشيف مضغوط مقدم كدليل")
            elif file_ext in ['.txt', '.md', '.json'] or file_ext in _CODE_EXTENSIONS:
                return DocumentProcessor._extract_from_text(file_path)
            elif file_ext in _IMAGE_EXTENSIONS:
                return DocumentProcessor._describe_image_file(file_path)
            else:
                # NEVER raise — unknown file types are recorded as evidence with
                # their metadata so the AI grader can still account for them.
                return DocumentProcessor._describe_binary_file(
                    file_path, f"Unknown file type ({file_ext})",
                    f"ملف بصيغة غير معروفة ({file_ext}) مقدم ضمن أعمال الطالب")

        except Exception as e:
            print(f"❌ Error processing file {file_path}: {e}")
            raise Exception(f"Failed to process document: {str(e)}")

    @staticmethod
    def count_images(file_path: str) -> int:
        """Count embedded images/screenshots in a document"""
        file_ext = os.path.splitext(file_path)[1].lower()
        try:
            if file_ext in ['.docx', '.doc']:
                return DocumentProcessor._count_images_docx(file_path)
            elif file_ext == '.pdf':
                return DocumentProcessor._count_images_pdf(file_path)
            elif file_ext == '.pptx':
                return DocumentProcessor._count_images_pptx(file_path)
        except Exception as e:
            print(f"⚠️ Image counting failed: {e}")
        return 0

    @staticmethod
    def _count_images_pptx(file_path: str) -> int:
        """Count picture shapes in a PPTX presentation."""
        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation(file_path)
            count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # PICTURE
                        count += 1
            return count
        except Exception:
            return 0

    @staticmethod
    def _count_images_docx(file_path: str) -> int:
        """Count images embedded in DOCX file"""
        doc = Document(file_path)
        count = 0
        # Count inline images (blipFill = embedded image)
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                count += 1
        return count

    @staticmethod
    def _count_images_pdf(file_path: str) -> int:
        """Count images in PDF pages"""
        count = 0
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                if page.images:
                    count += len(page.images)
        return count

    @staticmethod
    def extract_images(file_path: str, max_images: int = 20) -> List[Tuple[bytes, str]]:
        """
        Extract embedded images from a document.

        Args:
            file_path: Path to the document
            max_images: Maximum number of images to extract (0 = no limit)

        Returns:
            List of (image_bytes, mime_type) tuples
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        _IMAGE_EXTENSIONS = ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.tiff')
        try:
            if file_ext in ['.docx', '.doc']:
                return DocumentProcessor._extract_images_docx(file_path, max_images)
            elif file_ext == '.pdf':
                return DocumentProcessor._extract_images_pdf(file_path, max_images)
            elif file_ext == '.pptx':
                return DocumentProcessor._extract_images_pptx(file_path, max_images)
            elif file_ext in _IMAGE_EXTENSIONS:
                with open(file_path, 'rb') as f:
                    img_bytes = f.read()
                mime_map = {'.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', '.gif': 'image/gif', '.bmp': 'image/bmp', '.webp': 'image/webp', '.tiff': 'image/tiff'}
                return [(img_bytes, mime_map.get(file_ext, 'image/png'))]
        except Exception as e:
            print(f"⚠️ Image extraction failed: {e}")
        return []

    @staticmethod
    def _extract_images_docx(file_path: str, max_images: int) -> List[Tuple[bytes, str]]:
        """Extract embedded images from DOCX file (skips EMF/WMF for vision APIs)."""
        doc = Document(file_path)
        raw: List[Tuple[bytes, str]] = []
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_part = rel.target_part
                    image_bytes = image_part.blob
                    content_type = image_part.content_type  # e.g. "image/png"
                    if len(image_bytes) < 2048:
                        continue
                    raw.append((image_bytes, content_type))
                except Exception as e:
                    print(f"⚠️ Failed to extract image from DOCX rel: {e}")
                    continue
        images = filter_vision_images(raw)
        if max_images > 0:
            return images[:max_images]
        return images

    @staticmethod
    def _extract_images_pptx(file_path: str, max_images: int) -> List[Tuple[bytes, str]]:
        """Extract embedded images from PPTX slides (skips icons < 2KB)."""
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return []
        images: List[Tuple[bytes, str]] = []
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            if len(image_bytes) < 2048:
                                continue
                            content_type = image.content_type  # e.g. "image/png"
                            images.append((image_bytes, content_type))
                        except Exception as e:
                            print(f"⚠️ PPTX image extract failed for one shape: {e}")
                            continue
        except Exception as e:
            print(f"⚠️ PPTX image extraction failed: {e}")
        images = filter_vision_images(images)
        if max_images > 0:
            return images[:max_images]
        return images

    @staticmethod
    def _extract_images_pdf(file_path: str, max_images: int) -> List[Tuple[bytes, str]]:
        """Extract embedded images from PDF file"""
        import fitz  # PyMuPDF  # type: ignore

        images = []
        try:
            pdf_doc = fitz.open(file_path)
            for page_num in range(len(pdf_doc)):
                page = pdf_doc[page_num]
                image_list = page.get_images(full=True)
                for img_info in image_list:
                    xref = img_info[0]
                    try:
                        base_image = pdf_doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        ext = base_image["ext"]  # e.g. "png", "jpeg"
                        # Skip very small images (< 2KB)
                        if len(image_bytes) < 2048:
                            continue
                        mime_type = f"image/{ext}" if ext != "jpg" else "image/jpeg"
                        images.append((image_bytes, mime_type))
                        if max_images > 0 and len(images) >= max_images:
                            pdf_doc.close()
                            return images
                    except Exception:
                        continue
            pdf_doc.close()
        except ImportError:
            print("⚠️ PyMuPDF (fitz) not installed. PDF image extraction unavailable.")
        except Exception as e:
            print(f"⚠️ PDF image extraction failed: {e}")
        return images

    @staticmethod
    def _describe_binary_file(file_path: str, label: str, ar_note: str) -> str:
        """Marker string for a binary file we don't parse."""
        try:
            size_bytes = os.path.getsize(file_path)
            size_str = (f"{size_bytes / (1024 * 1024):.1f} MB" if size_bytes >= 1024 * 1024
                        else (f"{size_bytes / 1024:.1f} KB" if size_bytes >= 1024
                              else f"{size_bytes} bytes"))
        except Exception:
            size_str = "unknown size"
        fname = os.path.basename(file_path)
        return (f"[BINARY_FILE: {label}]\n"
                f"اسم الملف: {fname}\n"
                f"الحجم: {size_str}\n"
                f"ملاحظة: {ar_note}")

    @staticmethod
    def _describe_image_file(file_path: str) -> str:
        try:
            size_bytes = os.path.getsize(file_path)
            size_str = (f"{size_bytes / 1024:.1f} KB" if size_bytes >= 1024
                        else f"{size_bytes} bytes")
        except Exception:
            size_str = "unknown size"
        fname = os.path.basename(file_path)
        dims = "unknown dimensions"
        try:
            from PIL import Image as _Image  # type: ignore
            with _Image.open(file_path) as _img:
                dims = f"{_img.width}x{_img.height}"
        except Exception:
            pass
        return (f"[IMAGE_FILE: {fname}]\n"
                f"الأبعاد: {dims}\n"
                f"الحجم: {size_str}\n"
                f"ملاحظة: صورة مقدمة كدليل بصري")

    @staticmethod
    def _extract_from_pdf(file_path: str) -> str:
        """Extract text from PDF using pdfplumber"""
        text_content = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    # Add page marker for context
                    text_content.append(f"--- Page {i + 1} ---\n{text}\n")
        return "\n".join(text_content)

    @staticmethod
    def extract_text_with_image_count(
        file_path: str,
        max_chars: Optional[int] = None,
    ) -> Tuple[str, int]:
        """Single-pass DOCX extract + embedded image count (avoids opening file twice)."""
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext in (".docx", ".doc"):
            return DocumentProcessor._extract_from_docx_with_meta(file_path, max_chars=max_chars)
        return DocumentProcessor.extract_text(file_path, max_chars=max_chars), 0

    @staticmethod
    def _extract_from_docx(file_path: str, max_chars: Optional[int] = None) -> str:
        text, _ = DocumentProcessor._extract_from_docx_with_meta(file_path, max_chars=max_chars)
        return text

    @staticmethod
    def _extract_from_docx_with_meta(
        file_path: str,
        max_chars: Optional[int] = None,
    ) -> Tuple[str, int]:
        """Extract text from DOCX; optional max_chars stops early when set by caller."""
        doc = Document(file_path)
        image_count = sum(1 for rel in doc.part.rels.values() if "image" in rel.reltype)
        parts: List[str] = []
        used = 0
        budget = max_chars if max_chars and max_chars > 0 else 0

        def _append(chunk: str) -> bool:
            nonlocal used
            chunk = (chunk or "").strip()
            if not chunk:
                return True
            if budget and used >= budget:
                return False
            if budget:
                room = budget - used
                if len(chunk) > room:
                    chunk = chunk[:room]
            parts.append(chunk)
            used += len(chunk)
            return not budget or used < budget

        for para in doc.paragraphs:
            if not _append(para.text):
                break

        if not budget or used < budget:
            for table in doc.tables:
                if budget and used >= budget:
                    break
                for row in table.rows:
                    if budget and used >= budget:
                        break
                    row_text = []
                    for cell in row.cells:
                        cell_text = "\n".join(p.text for p in cell.paragraphs if p.text)
                        if cell_text.strip():
                            row_text.append(cell_text.strip())
                    if row_text and not _append(" | ".join(row_text)):
                        break

        joined = "\n".join(parts)
        if len(joined) < 100 and (not budget or used < budget):
            try:
                xml_text: List[str] = []
                for element in doc.element.body.iter():
                    if element.tag.endswith("t") and element.text:
                        xml_text.append(element.text)
                        if budget and sum(len(x) for x in xml_text) >= budget:
                            break
                xml_joined = "\n".join(xml_text)
                if len(xml_joined) > len(joined):
                    if budget and len(xml_joined) > budget:
                        xml_joined = xml_joined[:budget]
                    joined = xml_joined
            except Exception as e:
                print(f"XML extraction failed: {e}")

        if budget and len(joined) > budget:
            joined = joined[:budget]
        return joined, image_count

    @staticmethod
    def _extract_from_text(file_path: str) -> str:
        """Extract text from plain text files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    @staticmethod
    def _extract_from_excel(file_path: str) -> str:
        """Extract data from Excel files (.xlsx/.xls) including all sheets"""
        try:
            from openpyxl import load_workbook  # type: ignore
            wb = load_workbook(file_path, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"\n=== Sheet: {sheet_name} ===")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else '' for c in row]
                    if any(c.strip() for c in cells):
                        text_parts.append(' | '.join(cells))
            wb.close()
            return '\n'.join(text_parts)
        except Exception as e:
            return f"[EXCEL_FILE:{os.path.basename(file_path)} - Error: {e}]"

    @staticmethod
    def _extract_from_svg(file_path: str) -> str:
        """Extract text content and metadata from SVG files"""
        try:
            import xml.etree.ElementTree as ET
            tree = ET.parse(file_path)
            root = tree.getroot()
            # Get dimensions
            width = root.attrib.get('width', 'N/A')
            height = root.attrib.get('height', 'N/A')
            viewbox = root.attrib.get('viewBox', 'N/A')
            # Extract all text elements
            texts = []
            for elem in root.iter():
                tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                if tag in ('text', 'tspan') and elem.text and elem.text.strip():
                    texts.append(elem.text.strip())
            file_size = os.path.getsize(file_path)
            result = f"[SVG_FILE:{os.path.basename(file_path)}]\n"
            result += f"Dimensions: width={width}, height={height}, viewBox={viewbox}\n"
            result += f"File size: {file_size / 1024:.1f} KB\n"  # noqa: E226
            if texts:
                result += f"Text content: {', '.join(texts)}\n"
            # Count shapes
            shape_tags = ('rect', 'circle', 'ellipse', 'line', 'polyline', 'polygon', 'path')
            shape_count = sum(1 for elem in root.iter() if elem.tag.split('}')[-1] in shape_tags)
            result += f"Shape elements: {shape_count}\n"
            return result
        except Exception as e:
            return f"[SVG_FILE:{os.path.basename(file_path)} - Metadata extraction failed: {e}]"

    @staticmethod
    def _extract_video_metadata(file_path: str) -> str:
        """Extract metadata from video files (.mp4, etc.)"""
        try:
            file_size = os.path.getsize(file_path)
            file_name = os.path.basename(file_path)
            result = f"[VIDEO_FILE:{file_name}]\n"
            result += f"File type: {os.path.splitext(file_path)[1].upper()}\n"
            result += f"File size: {file_size / 1024:.1f} KB ({file_size / (1024 * 1024):.2f} MB)\n"
            # Try to get duration using ffprobe if available
            try:
                import subprocess  # noqa: delayed import
                cmd = ['ffprobe', '-v', 'quiet', '-show_entries', 'format=duration,bit_rate',
                       '-show_entries', 'stream=width,height,codec_name,r_frame_rate',
                       '-of', 'json', file_path]
                proc = subprocess.run(cmd, capture_output=True, text=True, timeout=10)
                if proc.returncode == 0:
                    import json
                    info = json.loads(proc.stdout)
                    fmt = info.get('format', {})
                    if 'duration' in fmt:
                        dur = float(fmt['duration'])
                        result += f"Duration: {dur:.1f} seconds\n"
                    if 'bit_rate' in fmt:
                        result += f"Bit rate: {int(fmt['bit_rate']) // 1000} kbps\n"
                    for s in info.get('streams', []):
                        if s.get('width'):
                            result += f"Resolution: {s['width']}x{s['height']}\n"
                            result += f"Codec: {s.get('codec_name', 'N/A')}\n"
                            if s.get('r_frame_rate'):
                                result += f"Frame rate: {s['r_frame_rate']}\n"
                            break
            except (FileNotFoundError, Exception):
                result += "(ffprobe not available - basic metadata only)\n"
            return result
        except Exception as e:
            return f"[VIDEO_FILE:{os.path.basename(file_path)} - Error: {e}]"

    @staticmethod
    def _extract_from_pptx(file_path: str) -> str:
        """Extract text content from PowerPoint (.pptx) — slide-by-slide,
        including text frames, tables, speaker notes, and shape titles.
        Returns a marker string if extraction fails (e.g. for legacy .ppt format).
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.ppt':
            # python-pptx only supports the modern .pptx format, not the old .ppt.
            return DocumentProcessor._describe_binary_file(
                file_path, "PowerPoint Presentation (legacy .ppt)",
                "ملف PowerPoint قديم (.ppt) — يُقبل كدليل بصري، لكن استخراج المحتوى غير مدعوم تلقائياً. "
                "للحصول على تقييم أعمق يُنصح بحفظ الملف بصيغة .pptx الحديثة.")
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return DocumentProcessor._describe_binary_file(
                file_path, "PowerPoint Presentation",
                "مكتبة python-pptx غير مثبتة. ثبّتها عبر: pip install python-pptx")
        try:
            prs = Presentation(file_path)
            text_parts = [f"[PPTX_FILE: {os.path.basename(file_path)} — {len(prs.slides)} شريحة]"]

            for slide_idx, slide in enumerate(prs.slides, start=1):
                text_parts.append(f"\n=== Slide {slide_idx} ===")

                # Slide layout/title (for context)
                try:
                    if slide.slide_layout and slide.slide_layout.name:
                        text_parts.append(f"[Layout: {slide.slide_layout.name}]")
                except Exception:
                    pass

                for shape in slide.shapes:
                    # Standard text frames (titles, content placeholders, text boxes)
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in para.runs)
                            if line.strip():
                                text_parts.append(line)

                    # Tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text:
                                    row_text.append(cell_text)
                            if row_text:
                                text_parts.append(" | ".join(row_text))

                # Speaker notes — students often put their explanations here
                try:
                    if slide.has_notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            text_parts.append(f"[Speaker Notes]: {notes_text}")
                except Exception:
                    pass

            return "\n".join(text_parts)
        except Exception as e:
            return f"[PPTX_FILE: {os.path.basename(file_path)} - Error extracting: {e}]"

    @staticmethod
    def _extract_from_scratch(file_path: str) -> str:
        """Extract project summary from a Scratch .sb3 file (ZIP with project.json)."""
        import zipfile
        import json as _json
        try:
            with zipfile.ZipFile(file_path, 'r') as zf:
                if 'project.json' not in zf.namelist():
                    return "[SCRATCH_FILE: project.json not found in .sb3]"
                project = _json.loads(zf.read('project.json').decode('utf-8'))

            targets = project.get("targets", [])
            lines = ["[SCRATCH_PROJECT]"]
            for target in targets:
                is_stage = target.get("isStage", False)
                name = target.get("name", "?")
                label = "Stage" if is_stage else f"Sprite: {name}"
                costumes = len(target.get("costumes", []))
                sounds = len(target.get("sounds", []))
                blocks = target.get("blocks", {})
                block_count = len(blocks) if isinstance(blocks, dict) else 0
                variables = len(target.get("variables", {}))
                lists_count = len(target.get("lists", {}))

                lines.append(f"\n--- {label} ---")
                lines.append(f"  Costumes: {costumes}, Sounds: {sounds}, Blocks: {block_count}")
                if variables:
                    lines.append(f"  Variables: {variables}")
                if lists_count:
                    lines.append(f"  Lists: {lists_count}")

                # List unique opcodes
                if isinstance(blocks, dict):
                    opcodes = sorted(set(
                        b.get("opcode", "") for b in blocks.values()
                        if isinstance(b, dict) and b.get("opcode")
                    ))
                    if opcodes:
                        lines.append(f"  Block types: {', '.join(opcodes)}")

            return "\n".join(lines)
        except Exception as e:
            return f"[SCRATCH_FILE: Error reading {os.path.basename(file_path)}: {e}]"

    @staticmethod
    def process_student_folder(folder_path: str) -> dict:
        """
        Process all files in a student submission folder.
        Returns a dict with extracted text, file inventory, and image data.
        """
        if not os.path.isdir(folder_path):
            raise FileNotFoundError(f"Folder not found: {folder_path}")

        _SUPPORTED = ('.docx', '.doc', '.pdf', '.xlsx', '.xls', '.pptx', '.txt',
                      '.png', '.jpg', '.jpeg', '.svg', '.mp4', '.avi', '.mov',
                      '.cs', '.gml', '.sb3', '.json', '.yaml', '.yml',
                      '.html', '.css', '.js', '.py', '.java', '.csv')
        _IMAGE_FOR_VISION = ('.png', '.jpg', '.jpeg')

        result = {
            'file_inventory': [],
            'text_content': [],
            'image_files': [],
            'combined_text': '',
        }

        for root_dir, dirs, files in os.walk(folder_path):
            # Skip hidden dirs
            dirs[:] = [d for d in dirs if not d.startswith('.')]
            for fname in sorted(files):
                if fname.startswith('.') or fname == 'desktop.ini':
                    continue
                fpath = os.path.join(root_dir, fname)
                ext = os.path.splitext(fname)[1].lower()
                rel_path = os.path.relpath(fpath, folder_path)
                file_size = os.path.getsize(fpath)

                file_info = {
                    'name': fname,
                    'path': rel_path,
                    'extension': ext,
                    'size_kb': round(file_size / 1024, 1),
                }
                result['file_inventory'].append(file_info)

                # Extract text from supported files
                if ext in _SUPPORTED:
                    try:
                        text = DocumentProcessor.extract_text(fpath)
                        if text and text.strip():
                            sep = '=' * 60
                            header = f"\n{sep}\n📄 FILE: {rel_path} ({ext}, {file_info['size_kb']}KB)\n{sep}\n"
                            result['text_content'].append(header + text)
                    except Exception as e:
                        result['text_content'].append(
                            f"\n📄 FILE: {rel_path} - [Error reading: {e}]\n"
                        )

                # Collect images for vision analysis
                if ext in _IMAGE_FOR_VISION and file_size > 2048:
                    try:
                        with open(fpath, 'rb') as f:
                            img_bytes = f.read()
                        mime = 'image/png' if ext == '.png' else 'image/jpeg'
                        result['image_files'].append({
                            'name': rel_path,
                            'bytes': img_bytes,
                            'mime': mime,
                        })
                    except Exception:
                        pass

        # Build combined text
        inventory_text = "\n📋 FILE INVENTORY (جرد الملفات):\n"
        for fi in result['file_inventory']:
            inventory_text += f"  • {fi['path']} ({fi['extension']}, {fi['size_kb']}KB)\n"
        result['combined_text'] = inventory_text + '\n'.join(result['text_content'])
        return result
